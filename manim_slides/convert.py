import mimetypes
import os
import platform
import shutil
import subprocess
import tempfile
import textwrap
import warnings
import webbrowser
from base64 import b64encode
from collections import deque
from enum import Enum
from importlib import resources
from pathlib import Path
from typing import Any, Callable, Optional, Union

import av
import click
import pptx
import requests
from bs4 import BeautifulSoup
from click import Context, Parameter
from jinja2 import Template
from lxml import etree
from PIL import Image
from pydantic import (
    BaseModel,
    ConfigDict,
    Field,
    FilePath,
    GetCoreSchemaHandler,
    PositiveFloat,
    PositiveInt,
    ValidationError,
)
from pydantic_core import CoreSchema, core_schema
from pydantic_extra_types.color import Color
from tqdm import tqdm

from . import templates
from .commons import folder_path_option, verbosity_option
from .config import PresentationConfig
from .logger import logger
from .present import get_scenes_presentation_config


def open_with_default(file: Path) -> None:
    system = platform.system()
    if system == "Darwin":
        subprocess.call(("open", str(file)))
    elif system == "Windows":
        os.startfile(str(file))  # type: ignore[attr-defined]
    else:
        subprocess.call(("xdg-open", str(file)))


def validate_config_option(
    ctx: Context, param: Parameter, value: Any
) -> dict[str, str]:
    config = {}

    for c_option in value:
        try:
            key, value = c_option.split("=")
            config[key] = value
        except ValueError:
            raise click.BadParameter(
                f"Configuration options `{c_option}` could not be parsed into "
                "a proper (key, value) pair. "
                "Please use an `=` sign to separate key from value."
            ) from None

    return config


def file_to_data_uri(file: Path) -> str:
    """Read a video and return the corresponding data-uri."""
    b64 = b64encode(file.read_bytes()).decode("ascii")
    mime_type = mimetypes.guess_type(file)[0] or "video/mp4"

    return f"data:{mime_type};base64,{b64}"


def get_duration_ms(file: Path) -> float:
    """Read a video and return its duration in milliseconds."""
    with av.open(str(file)) as container:
        video = container.streams.video[0]

        return float(1000 * video.duration * video.time_base)


def read_image_from_video_file(file: Path, frame_index: "FrameIndex") -> Image:
    """Read a image from a video file at a given index."""
    with av.open(str(file)) as container:
        frames = container.decode(video=0)

        if frame_index == FrameIndex.last:
            (frame,) = deque(frames, 1)
        else:
            frame = next(frames)

        return frame.to_image()


class Converter(BaseModel):  # type: ignore
    presentation_configs: list[PresentationConfig]
    assets_dir: str = Field(
        "{basename}_assets",
        description="Assets folder.\nThis is a template string that accepts 'dirname', 'basename', and 'ext' as variables.\nThose variables are obtained from the output filename.",
    )
    template: Optional[Path] = Field(None, description="Custom template file to use.")

    def convert_to(self, dest: Path) -> None:
        """Convert self, i.e., a list of presentations, into a given format."""
        raise NotImplementedError

    def load_template(self) -> str:
        """
        Return the template as a string.

        An empty string is returned if no template is used.
        """
        return ""

    def open(self, file: Path) -> None:
        """Open a file, generated with converter, using appropriate application."""
        open_with_default(file)

    @classmethod
    def from_string(cls, s: str) -> type["Converter"]:
        """Return the appropriate converter from a string name."""
        return {
            "html": RevealJS,
            "pdf": PDF,
            "pptx": PowerPoint,
            "zip": HtmlZip,
        }[s]


class Str(str):
    """A simple string, but quoted when needed."""

    # This fixes pickling issue on Python 3.8
    __reduce_ex__ = str.__reduce_ex__

    @classmethod
    def __get_pydantic_core_schema__(
        cls, source_type: Any, handler: GetCoreSchemaHandler
    ) -> CoreSchema:
        return core_schema.str_schema()

    def __str__(self) -> str:
        """Ensure that the string is correctly quoted."""
        if self in ["true", "false", "null"]:
            return self
        else:
            return f"'{super().__str__()}'"


class StrEnum(Enum):
    def __str__(self) -> str:
        return str(self.value)


Function = str  # Basically, anything


class JsTrue(str, StrEnum):
    true = "true"


class JsFalse(str, StrEnum):
    false = "false"


class JsBool(Str, StrEnum):  # type: ignore
    true = "true"
    false = "false"


class JsNull(Str, StrEnum):  # type: ignore
    null = "null"


class ControlsLayout(Str, StrEnum):  # type: ignore
    edges = "edges"
    bottom_right = "bottom-right"


class ControlsBackArrows(Str, StrEnum):  # type: ignore
    faded = "faded"
    hidden = "hidden"
    visibly = "visibly"


class SlideNumber(Str, StrEnum):  # type: ignore
    true = "true"
    false = "false"
    hdotv = "h.v"
    handv = "h/v"
    c = "c"
    candt = "c/t"


class ShowSlideNumber(Str, StrEnum):  # type: ignore
    all = "all"
    print = "print"
    speaker = "speaker"


class KeyboardCondition(Str, StrEnum):  # type: ignore
    null = "null"
    focused = "focused"


class NavigationMode(Str, StrEnum):  # type: ignore
    default = "default"
    linear = "linear"
    grid = "grid"


class AutoPlayMedia(Str, StrEnum):  # type: ignore
    null = "null"
    true = "true"
    false = "false"


PreloadIframes = AutoPlayMedia


class AutoAnimateMatcher(Str, StrEnum):  # type: ignore
    null = "null"


class AutoAnimateEasing(Str, StrEnum):  # type: ignore
    ease = "ease"


AutoSlide = Union[PositiveInt, JsFalse]


class AutoSlideMethod(Str, StrEnum):  # type: ignore
    null = "null"


MouseWheel = Union[JsNull, float]


class Transition(Str, StrEnum):  # type: ignore
    none = "none"
    fade = "fade"
    slide = "slide"
    convex = "convex"
    concave = "concave"
    zoom = "zoom"


class TransitionSpeed(Str, StrEnum):  # type: ignore
    default = "default"
    fast = "fast"
    slow = "slow"


class BackgroundSize(Str, StrEnum):  # type: ignore
    # From: https://developer.mozilla.org/en-US/docs/Web/CSS/background-size
    # TODO: support more background size
    contain = "contain"
    cover = "cover"


BackgroundTransition = Transition


class Display(Str, StrEnum):  # type: ignore
    block = "block"


class RevealTheme(str, StrEnum):
    black = "black"
    white = "white"
    league = "league"
    beige = "beige"
    sky = "sky"
    night = "night"
    serif = "serif"
    simple = "simple"
    soralized = "solarized"
    blood = "blood"
    moon = "moon"
    black_contrast = "black-contrast"
    white_contrast = "white-contrast"
    dracula = "dracula"


class RevealJS(Converter):
    """
    RevealJS options.

    Please check out https://revealjs.com/config/ for more details.
    """

    # Export option:
    one_file: bool = Field(
        False, description="Embed all assets (e.g., animations) inside the HTML."
    )
    offline: bool = Field(
        False, description="Download remote assets for offline presentation."
    )
    # Presentation size options from RevealJS
    width: Union[Str, int] = Field(
        Str("100%"), description="Width of the presentation."
    )
    height: Union[Str, int] = Field(
        Str("100%"), description="Height of the presentation."
    )
    margin: float = Field(0.04, description="Margin to use around the content.")
    min_scale: float = Field(
        0.2, description="Bound for smallest possible scale to apply to content."
    )
    max_scale: float = Field(
        2.0, description="Bound for large possible scale to apply to content."
    )
    # Configuration options from RevealJS
    controls: JsBool = Field(
        JsBool.false, description="Display presentation control arrows."
    )
    controls_tutorial: JsBool = Field(
        JsBool.true, description="Help the user learn the controls by providing hints."
    )
    controls_layout: ControlsLayout = Field(
        ControlsLayout.bottom_right, description="Determine where controls appear."
    )
    controls_back_arrows: ControlsBackArrows = Field(
        ControlsBackArrows.faded,
        description="Visibility rule for backwards navigation arrows.",
    )
    progress: JsBool = Field(
        JsBool.false, description="Display a presentation progress bar."
    )
    slide_number: SlideNumber = Field(
        SlideNumber.false, description="Display the page number of the current slide."
    )
    show_slide_number: Union[ShowSlideNumber, Function] = Field(
        ShowSlideNumber.all,
        description="Can be used to limit the contexts in which the slide number appears.",
    )
    hash_one_based_index: JsBool = Field(
        JsBool.false,
        description="Use 1 based indexing for # links to match slide number (default is zero based).",
    )
    hash: JsBool = Field(
        JsBool.false,
        description="Add the current slide number to the URL hash so that reloading the page/copying the URL will return you to the same slide.",
    )
    respond_to_hash_changes: JsBool = Field(
        JsBool.false,
        description="Flags if we should monitor the hash and change slides accordingly.",
    )
    jump_to_slide: JsBool = Field(
        JsBool.true,
        description="Enable support for jump-to-slide navigation shortcuts.",
    )
    history: JsBool = Field(
        JsBool.false,
        description="Push each slide change to the browser history.  Implies `hash: true`.",
    )
    keyboard: JsBool = Field(
        JsBool.true, description="Enable keyboard shortcuts for navigation."
    )
    keyboard_condition: Union[KeyboardCondition, Function] = Field(
        KeyboardCondition.null,
        description="Optional function that blocks keyboard events when retuning false.",
    )
    disable_layout: JsBool = Field(
        JsBool.false,
        description="Disable the default reveal.js slide layout (scaling and centering) so that you can use custom CSS layout.",
    )
    overview: JsBool = Field(JsBool.true, description="Enable the slide overview mode.")
    center: JsBool = Field(JsBool.true, description="Vertical centering of slides.")
    touch: JsBool = Field(
        JsBool.true, description="Enable touch navigation on devices with touch input."
    )
    loop: JsBool = Field(JsBool.false, description="Loop the presentation.")
    rtl: JsBool = Field(
        JsBool.false, description="Change the presentation direction to be RTL."
    )
    navigation_mode: NavigationMode = Field(
        NavigationMode.default,
        description="Change the behavior of our navigation directions.",
    )
    shuffle: JsBool = Field(
        JsBool.false,
        description="Randomize the order of slides each time the presentation loads.",
    )
    fragments: JsBool = Field(
        JsBool.true, description="Turns fragment on and off globally."
    )
    fragment_in_url: JsBool = Field(
        JsBool.true,
        description="Flag whether to include the current fragment in the URL, so that reloading brings you to the same fragment position.",
    )
    embedded: JsBool = Field(
        JsBool.false,
        description="Flag if the presentation is running in an embedded mode, i.e. contained within a limited portion of the screen.",
    )
    help: JsBool = Field(
        JsBool.true,
        description="Flag if we should show a help overlay when the question-mark key is pressed.",
    )
    pause: JsBool = Field(
        JsBool.true,
        description="Flag if it should be possible to pause the presentation (blackout).",
    )
    show_notes: JsBool = Field(
        JsBool.false,
        description="Flag if speaker notes should be visible to all viewers.",
    )
    auto_play_media: AutoPlayMedia = Field(
        AutoPlayMedia.null,
        description="Global override for autolaying embedded media (video/audio/iframe).",
    )
    preload_iframes: PreloadIframes = Field(
        PreloadIframes.null,
        description="Global override for preloading lazy-loaded iframes.",
    )
    auto_animate: JsBool = Field(
        JsBool.true, description="Can be used to globally disable auto-animation."
    )
    auto_animate_matcher: Union[AutoAnimateMatcher, Function] = Field(
        AutoAnimateMatcher.null,
        description="Optionally provide a custom element matcher that will be used to dictate which elements we can animate between.",
    )
    auto_animate_easing: AutoAnimateEasing = Field(
        AutoAnimateEasing.ease,
        description="Default settings for our auto-animate transitions, can be overridden per-slide or per-element via data arguments.",
    )
    auto_animate_duration: float = Field(
        1.0, description="See 'auto_animate_easing' documentation."
    )
    auto_animate_unmatched: JsBool = Field(
        JsBool.true, description="See 'auto_animate_easing' documentation."
    )
    auto_animate_styles: list[str] = Field(
        default_factory=lambda: [
            "opacity",
            "color",
            "background-color",
            "padding",
            "font-size",
            "line-height",
            "letter-spacing",
            "border-width",
            "border-color",
            "border-radius",
            "outline",
            "outline-offset",
        ],
        description="CSS properties that can be auto-animated.",
    )
    auto_slide: AutoSlide = Field(
        0, description="Control automatic progression to the next slide."
    )
    auto_slide_stoppable: JsBool = Field(
        JsBool.true, description="Stop auto-sliding after user input."
    )
    auto_slide_method: Union[AutoSlideMethod, Function] = Field(
        AutoSlideMethod.null,
        description="Use this method for navigation when auto-sliding (defaults to navigateNext).",
    )
    default_timing: Union[JsNull, int] = Field(
        JsNull.null,
        description="Specify the average time in seconds that you think you will spend presenting each slide.",
    )
    mouse_wheel: JsBool = Field(
        JsBool.false, description="Enable slide navigation via mouse wheel."
    )
    preview_links: JsBool = Field(
        JsBool.false, description="Open links in an iframe preview overlay."
    )
    post_message: JsBool = Field(
        JsBool.true, description="Expose the reveal.js API through window.postMessage."
    )
    post_message_events: JsBool = Field(
        JsBool.false,
        description="Dispatch all reveal.js events to the parent window through postMessage.",
    )
    focus_body_on_page_visibility_change: JsBool = Field(
        JsBool.true,
        description="Focus body when page changes visibility to ensure keyboard shortcuts work.",
    )
    transition: Transition = Field(Transition.none, description="Transition style.")
    transition_speed: TransitionSpeed = Field(
        TransitionSpeed.default, description="Transition speed."
    )
    background_size: BackgroundSize = Field(
        BackgroundSize.contain, description="Background size attribute for each video."
    )  # Not in RevealJS
    background_transition: BackgroundTransition = Field(
        BackgroundTransition.none,
        description="Transition style for full page slide backgrounds.",
    )
    pdf_max_pages_per_slide: Union[int, str] = Field(
        "Number.POSITIVE_INFINITY",
        description="The maximum number of pages a single slide can expand onto when printing to PDF, unlimited by default.",
    )
    pdf_separate_fragments: JsBool = Field(
        JsBool.true, description="Print each fragment on a separate slide."
    )
    pdf_page_height_offset: int = Field(
        -1,
        description="Offset used to reduce the height of content within exported PDF pages.",
    )
    view_distance: int = Field(
        3, description="Number of slides away from the current that are visible."
    )
    mobile_view_distance: int = Field(
        2,
        description="Number of slides away from the current that are visible on mobile devices.",
    )
    display: Display = Field(
        Display.block, description="The display mode that will be used to show slides."
    )
    hide_inactive_cursor: JsBool = Field(
        JsBool.true, description="Hide cursor if inactive."
    )
    hide_cursor_time: int = Field(
        5000, description="Time before the cursor is hidden (in ms)."
    )
    # Appearance options from RevealJS
    background_color: Color = Field(
        "black",
        description="Background color used in slides, not relevant if videos fill the whole area.",
    )
    reveal_version: str = Field("5.1.0", description="RevealJS version.")
    reveal_theme: RevealTheme = Field(
        RevealTheme.black, description="RevealJS version."
    )
    title: str = Field("Manim Slides", description="Presentation title.")
    # Pydantic options
    model_config = ConfigDict(use_enum_values=True, extra="forbid")

    def load_template(self) -> str:
        """Return the RevealJS HTML template as a string."""
        if isinstance(self.template, Path):
            return self.template.read_text()

        return resources.files(templates).joinpath("revealjs.html").read_text()

    def open(self, file: Path) -> None:
        webbrowser.open(file.absolute().as_uri())

    def convert_to(self, dest: Path) -> None:  # noqa: C901
        """
        Convert this configuration into a RevealJS HTML presentation, saved to
        DEST.
        """
        dirname = dest.parent
        basename = dest.stem
        ext = dest.suffix

        assets_dir = Path(
            self.assets_dir.format(dirname=dirname, basename=basename, ext=ext)
        )
        full_assets_dir = dirname / assets_dir

        if not self.one_file or self.offline:
            logger.debug(f"Assets will be saved to: {full_assets_dir}")

        if not self.one_file:
            num_presentation_configs = len(self.presentation_configs)

            if num_presentation_configs > 1:
                # Prevent possible name collision, see:
                # https://github.com/jeertmans/manim-slides/issues/428
                # With ManimCE, this can happen when caching is disabled as filenames are
                #   'uncached_000x.mp4'
                # With ManimGL, this can easily occur since filenames are just basic integers...
                num_digits = len(str(num_presentation_configs - 1))

                def prefix(i: int) -> str:
                    return f"s{i:0{num_digits}d}_"

            else:

                def prefix(i: int) -> str:
                    return ""

            full_assets_dir.mkdir(parents=True, exist_ok=True)
            for i, presentation_config in enumerate(self.presentation_configs):
                presentation_config.copy_to(
                    full_assets_dir, include_reversed=False, prefix=prefix(i)
                )

        dest.parent.mkdir(parents=True, exist_ok=True)

        with open(dest, "w") as f:
            revealjs_template = Template(self.load_template())

            options = self.model_dump()

            if assets_dir is not None:
                options["assets_dir"] = assets_dir

            has_notes = any(
                slide_config.notes != ""
                for presentation_config in self.presentation_configs
                for slide_config in presentation_config.slides
            )

            content = revealjs_template.render(
                file_to_data_uri=file_to_data_uri,
                get_duration_ms=get_duration_ms,
                has_notes=has_notes,
                env=os.environ,
                prefix=prefix if not self.one_file else None,
                **options,
            )
            # If not offline, write the content to the file
            if not self.offline:
                f.write(content)
                return

            # If offline, download remote assets and store them in the assets folder
            soup = BeautifulSoup(content, "html.parser")
            session = requests.Session()

            for tag, inner in [("link", "href"), ("script", "src")]:
                for item in soup.find_all(tag):
                    if item.has_attr(inner) and (link := item[inner]).startswith(
                        "http"
                    ):
                        asset_name = link.rsplit("/", 1)[1]
                        asset = session.get(link)
                        if self.one_file:
                            # If it is a CSS file, inline it
                            if tag == "link" and "stylesheet" in item["rel"]:
                                item.decompose()
                                style = soup.new_tag("style")
                                style.string = asset.text
                                soup.head.append(style)
                            # If it is a JS file, inline it
                            elif tag == "script":
                                item.decompose()
                                script = soup.new_tag("script")
                                script.string = asset.text
                                soup.head.append(script)
                            else:
                                raise ValueError(
                                    f"Unable to inline {tag} asset: {link}"
                                )
                        else:
                            full_assets_dir.mkdir(parents=True, exist_ok=True)
                            with open(full_assets_dir / asset_name, "wb") as asset_file:
                                asset_file.write(asset.content)

                            item[inner] = str(assets_dir / asset_name)

            content = str(soup)
            f.write(content)


class HtmlZip(RevealJS):
    def open(self, file: Path) -> None:
        super(RevealJS, self).open(file)  # Override opening with web browser

    def convert_to(self, dest: Path) -> None:
        """
        Convert this configuration into a zipped RevealJS HTML presentation, saved to
        DEST.
        """
        with tempfile.TemporaryDirectory() as directory_name:
            directory = Path(directory_name)

            html_file = directory / dest.with_suffix(".html").name

            super().convert_to(html_file)
            shutil.make_archive(str(dest.with_suffix("")), "zip", directory_name)


class FrameIndex(str, Enum):
    first = "first"
    last = "last"

    def __repr__(self) -> str:
        return self.value


class PDF(Converter):
    frame_index: FrameIndex = Field(
        FrameIndex.last,
        description="What frame (first or last) is used to represent each slide.",
    )
    resolution: PositiveFloat = Field(
        100.0, description="Image resolution use for saving frames."
    )
    model_config = ConfigDict(use_enum_values=True, extra="forbid")

    def convert_to(self, dest: Path) -> None:
        """Convert this configuration into a PDF presentation, saved to DEST."""
        images = []

        for i, presentation_config in enumerate(self.presentation_configs):
            for slide_config in tqdm(
                presentation_config.slides,
                desc=f"Generating video slides for config {i + 1}",
                leave=False,
            ):
                images.append(
                    read_image_from_video_file(slide_config.file, self.frame_index)
                )

        dest.parent.mkdir(parents=True, exist_ok=True)

        images[0].save(
            dest,
            "PDF",
            resolution=self.resolution,
            save_all=True,
            append_images=images[1:],
        )


class PowerPoint(Converter):
    left: PositiveInt = Field(
        0, description="Horizontal offset where the video is placed from left border."
    )
    top: PositiveInt = Field(
        0, description="Vertical offset where the video is placed from top border."
    )
    width: PositiveInt = Field(
        1280,
        description="Width of the slides.\nThis should match the resolution of the presentation.",
    )
    height: PositiveInt = Field(
        720,
        description="Height of the slides.\nThis should match the resolution of the presentation.",
    )
    auto_play_media: bool = Field(
        True, description="Automatically play animations when changing slide."
    )
    poster_frame_image: Optional[FilePath] = Field(
        None,
        description="Optional image to use when animations are not playing.\n"
        "By default, the first frame of each animation is used.\nThis is important to avoid blinking effects between slides.",
    )
    model_config = ConfigDict(use_enum_values=True, extra="forbid")

    def convert_to(self, dest: Path) -> None:
        """Convert this configuration into a PowerPoint presentation, saved to DEST."""
        prs = pptx.Presentation()
        prs.slide_width = self.width * 9525
        prs.slide_height = self.height * 9525

        layout = prs.slide_layouts[6]  # Should be blank

        # From GitHub issue comment:
        # - https://github.com/scanny/python-pptx/issues/427#issuecomment-856724440
        def auto_play_media(
            media: pptx.shapes.picture.Movie, loop: bool = False
        ) -> None:
            el_id = xpath(media.element, ".//p:cNvPr")[0].attrib["id"]
            el_cnt = xpath(
                media.element.getparent().getparent().getparent(),
                f'.//p:timing//p:video//p:spTgt[@spid="{el_id}"]',
            )[0]
            cond = xpath(el_cnt.getparent().getparent(), ".//p:cond")[0]
            cond.set("delay", "0")

            if loop:
                ctn = xpath(el_cnt.getparent().getparent(), ".//p:cTn")[0]
                ctn.set("repeatCount", "indefinite")

        def xpath(el: etree.Element, query: str) -> etree.XPath:
            nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
            return etree.ElementBase.xpath(el, query, namespaces=nsmap)

        with tempfile.TemporaryDirectory() as directory_name:
            directory = Path(directory_name)
            frame_number = 0
            for i, presentation_config in enumerate(self.presentation_configs):
                for slide_config in tqdm(
                    presentation_config.slides,
                    desc=f"Generating video slides for config {i + 1}",
                    leave=False,
                ):
                    file = slide_config.file

                    mime_type = mimetypes.guess_type(file)[0]

                    if self.poster_frame_image is None:
                        poster_frame_image = str(directory / f"{frame_number}.png")
                        image = read_image_from_video_file(
                            file, frame_index=FrameIndex.first
                        )
                        image.save(poster_frame_image)

                        frame_number += 1
                    else:
                        poster_frame_image = str(self.poster_frame_image)

                    slide = prs.slides.add_slide(layout)
                    movie = slide.shapes.add_movie(
                        str(file),
                        self.left,
                        self.top,
                        self.width * 9525,
                        self.height * 9525,
                        poster_frame_image=poster_frame_image,
                        mime_type=mime_type,
                    )
                    if slide_config.notes != "":
                        slide.notes_slide.notes_text_frame.text = slide_config.notes

                    if self.auto_play_media:
                        auto_play_media(movie, loop=slide_config.loop)

            dest.parent.mkdir(parents=True, exist_ok=True)
            prs.save(dest)


def show_config_options(function: Callable[..., Any]) -> Callable[..., Any]:
    """Wrap a function to add a '--show-config' option."""

    def callback(ctx: Context, _param: Parameter, value: bool) -> None:
        if not value or ctx.resilient_parsing:
            return

        if "to" in ctx.params:
            to = ctx.params["to"]
            cls = Converter.from_string(to)
        elif "dest" in ctx.params:
            dest = Path(ctx.params["dest"])
            fmt = dest.suffix[1:].lower()
            try:
                cls = Converter.from_string(fmt)
            except KeyError:
                logger.warning(
                    f"Could not guess conversion format from {dest!s}, defaulting to HTML."
                )
                cls = RevealJS
        else:
            cls = RevealJS

        if doc := getattr(cls, "__doc__", ""):
            click.echo(textwrap.dedent(doc))

        for key, field in cls.model_fields.items():
            if field.is_required():
                continue

            default = field.get_default(call_default_factory=True)
            click.echo(click.style(key, bold=True) + f": {default}")
            if description := field.description:
                click.secho(textwrap.indent(description, prefix="# "), dim=True)

        ctx.exit()

    return click.option(  # type: ignore
        "--show-config",
        is_flag=True,
        help="Show supported options for given format and exit.",
        default=None,
        expose_value=False,
        show_envvar=True,
        callback=callback,
    )(function)


def show_template_option(function: Callable[..., Any]) -> Callable[..., Any]:
    """Wrap a function to add a '--show-template' option."""

    def callback(ctx: Context, param: Parameter, value: bool) -> None:
        if not value or ctx.resilient_parsing:
            return

        if "to" in ctx.params:
            to = ctx.params["to"]
            cls = Converter.from_string(to)
        elif "dest" in ctx.params:
            dest = Path(ctx.params["dest"])
            fmt = dest.suffix[1:].lower()
            try:
                cls = Converter.from_string(fmt)
            except KeyError:
                logger.warning(
                    f"Could not guess conversion format from {dest!s}, defaulting to HTML."
                )
                cls = RevealJS
        else:
            cls = RevealJS

        template = ctx.params.get("template", None)

        converter = cls(presentation_configs=[], template=template)
        click.echo(converter.load_template())

        ctx.exit()

    return click.option(  # type: ignore
        "--show-template",
        is_flag=True,
        help="Show the template (currently) used for a given conversion format and exit.",
        default=None,
        expose_value=False,
        show_envvar=True,
        callback=callback,
    )(function)


@click.command()
@click.argument("scenes", nargs=-1)
@folder_path_option
@click.argument("dest", type=click.Path(dir_okay=False, path_type=Path))
@click.option(
    "--to",
    type=click.Choice(["auto", "html", "pdf", "pptx", "zip"], case_sensitive=False),
    metavar="FORMAT",
    default="auto",
    show_default=True,
    help="Set the conversion format to use. Use 'auto' to detect format from DEST.",
)
@click.option(
    "--open",
    "open_result",
    is_flag=True,
    help="Open the newly created file using the appropriate application.",
)
@click.option(
    "-c",
    "--config",
    "config_options",
    multiple=True,
    callback=validate_config_option,
    help="Configuration options passed to the converter. "
    "E.g., pass '-cslide_number=true' to display slide numbers.",
)
@click.option(
    "--use-template",
    "template",
    metavar="FILE",
    type=click.Path(exists=True, dir_okay=False, path_type=Path),
    help="Use the template given by FILE instead of default one. "
    "To echo the default template, use '--show-template'.",
)
@click.option(
    "--one-file",
    is_flag=True,
    help="Embed all local assets (e.g., video files) in the output file. "
    "The is a convenient alias to '-cone_file=true'.",
)
@click.option(
    "--offline",
    is_flag=True,
    help="Download any remote content and store it in the assets folder. "
    "The is a convenient alias to '-coffline=true'.",
)
@show_template_option
@show_config_options
@verbosity_option
def convert(
    scenes: list[str],
    folder: Path,
    dest: Path,
    to: str,
    open_result: bool,
    config_options: dict[str, str],
    template: Optional[Path],
    offline: bool,
    one_file: bool,
) -> None:
    """Convert SCENE(s) into a given format and writes the result in DEST."""
    presentation_configs = get_scenes_presentation_config(scenes, folder)

    try:
        if to == "auto":
            fmt = dest.suffix[1:].lower()
            try:
                cls = Converter.from_string(fmt)
            except KeyError:
                logger.warning(
                    f"Could not guess conversion format from {dest!s}, defaulting to HTML."
                )
                cls = RevealJS
        else:
            cls = Converter.from_string(to)

        if (
            one_file
            and issubclass(cls, (RevealJS, HtmlZip))
            and "one_file" not in config_options
        ):
            config_options["one_file"] = "true"

        # Change data_uri to one_file and print a warning if present
        if "data_uri" in config_options:
            warnings.warn(
                "The 'data_uri' configuration option is deprecated and will be "
                "removed in the next major version. "
                "Use 'one_file' instead.",
                DeprecationWarning,
                stacklevel=2,
            )
            config_options["one_file"] = (
                config_options["one_file"]
                if "one_file" in config_options
                else config_options.pop("data_uri")
            )

        if (
            offline
            and issubclass(cls, (RevealJS, HtmlZip))
            and "offline" not in config_options
        ):
            config_options["offline"] = "true"

        converter = cls(
            presentation_configs=presentation_configs,
            template=template,
            **config_options,
        )

        converter.convert_to(dest)

        if open_result:
            converter.open(dest)

    except ValidationError as e:
        errors = e.errors()

        msg = [
            f"{len(errors)} error(s) occurred with configuration options for '{to}', see below."
        ]

        for error in errors:
            option = error["loc"][0]
            _msg = error["msg"]
            msg.append(f"Option '{option}': {_msg}")

        raise click.UsageError("\n".join(msg)) from None
